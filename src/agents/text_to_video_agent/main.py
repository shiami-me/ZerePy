from typing import List, Optional, Tuple, Dict, Any
from pydantic import BaseModel, Field
from pathlib import Path
import os
import shutil
import logging
import io
import requests
from datetime import datetime

from PIL import Image, ImageDraw, ImageFont, ImageEnhance
import cv2
import numpy as np
import cairocffi as cairo

from langchain_core.tools import tool
from langchain_google_genai import ChatGoogleGenerativeAI
from langgraph.prebuilt import create_react_agent
from langgraph_supervisor import create_supervisor

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
def setup_logging() -> logging.Logger:
    """Set up logging configuration."""
    logs_dir = "logs"
    os.makedirs(logs_dir, exist_ok=True)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    log_file = os.path.join(logs_dir, f"presentation_creation_{timestamp}.log")
    
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler(log_file),
            logging.StreamHandler()
        ]
    )
    return logging.getLogger(__name__)

logger = setup_logging()

# Data Models

class SlideTheme(BaseModel):
    """Model for slide theme configuration."""
    background_color: str = Field(default="#FFFFFF")
    title_font: str = Field(default="Arial")
    body_font: str = Field(default="Arial")
    title_color: str = Field(default="#000000")
    body_color: str = Field(default="#333333")
    accent_color: str = Field(default="#0066CC")
    padding: int = Field(default=50)

class ImageGenerationConfig(BaseModel):
    """Configuration for image generation."""
    width: int = Field(default=1920)
    height: int = Field(default=1080)
    dpi: int = Field(default=300)
    quality: int = Field(default=95)
    enable_animations: bool = Field(default=True)
    transition_frames: int = Field(default=10)

class SlideElement:
    """Base class for slide elements."""
    def __init__(self, x: int, y: int, width: int, height: int) -> None:
        self.x = x
        self.y = y
        self.width = width
        self.height = height
    
    def draw(self, context: cairo.Context) -> None:
        raise NotImplementedError

class TextElement(SlideElement):
    """Handles text elements with formatting."""
    def __init__(self, x: int, y: int, width: int, height: int, 
                 text: str, font_size: int, font_family: str,
                 color: str, bold: bool = False, italic: bool = False) -> None:
        super().__init__(x, y, width, height)
        self.text = text
        self.font_size = font_size
        self.font_family = font_family
        self.color = color
        self.bold = bold
        self.italic = italic
    
    def draw(self, context: cairo.Context) -> None:
        context.save()
        
        # Set font properties
        weight = cairo.FONT_WEIGHT_BOLD if self.bold else cairo.FONT_WEIGHT_NORMAL
        slant = cairo.FONT_SLANT_ITALIC if self.italic else cairo.FONT_SLANT_NORMAL
        context.select_font_face(self.font_family, slant, weight)
        context.set_font_size(self.font_size)
        
        # Set color (assumes color string like "#RRGGBB")
        r, g, b = [int(self.color[i:i+2], 16)/255 for i in (1, 3, 5)]
        context.set_source_rgb(r, g, b)
        
        # Draw text with simple word wrap
        words = self.text.split()
        line = []
        y_offset = 0
        
        for word in words:
            line.append(word)
            test_line = ' '.join(line)
            extents = context.text_extents(test_line)
            
            if extents.width > self.width:
                line.pop()  # remove last word and draw the current line
                final_line = ' '.join(line)
                context.move_to(self.x, self.y + y_offset)
                context.show_text(final_line)
                line = [word]
                y_offset += self.font_size * 1.5
            
        # Draw the last line
        final_line = ' '.join(line)
        context.move_to(self.x, self.y + y_offset)
        context.show_text(final_line)
        
        context.restore()

class ShapeElement(SlideElement):
    """Handles geometric shapes."""
    def __init__(self, x: int, y: int, width: int, height: int,
                 shape_type: str, fill_color: str, border_color: Optional[str] = None,
                 border_width: float = 1.0) -> None:
        super().__init__(x, y, width, height)
        self.shape_type = shape_type
        self.fill_color = fill_color
        self.border_color = border_color
        self.border_width = border_width
    
    def draw(self, context: cairo.Context) -> None:
        context.save()
        
        # Set fill color
        r, g, b = [int(self.fill_color[i:i+2], 16)/255 for i in (1, 3, 5)]
        context.set_source_rgb(r, g, b)
        
        if self.shape_type == "rectangle":
            context.rectangle(self.x, self.y, self.width, self.height)
        elif self.shape_type == "ellipse":
            context.save()
            # Center and scale the ellipse
            context.translate(self.x + self.width / 2, self.y + self.height / 2)
            context.scale(self.width / 2, self.height / 2)
            context.arc(0, 0, 1, 0, 2 * np.pi)
            context.restore()
        
        context.fill()
        
        if self.border_color:
            context.set_line_width(self.border_width)
            r, g, b = [int(self.border_color[i:i+2], 16)/255 for i in (1, 3, 5)]
            context.set_source_rgb(r, g, b)
            context.stroke()
        
        context.restore()

class SlideData(BaseModel):
    """Model for individual slide data."""
    topic: str = Field(description="The topic/title of the slide")
    content: str = Field(description="The content/body of the slide")
    bullet_points: Optional[List[str]] = Field(default=None, description="Optional bullet points for the slide")

class PresentationRequest(BaseModel):
    """Model for presentation request."""
    num_slides: int = Field(description="Number of slides requested")
    main_topic: str = Field(description="Main topic of the presentation")
    target_audience: Optional[str] = Field(default="general", description="Target audience for the presentation")
    style: Optional[str] = Field(default="professional", description="Style of the presentation")

class SlideContent(BaseModel):
    """Model for complete slide content."""
    slides: List[SlideData] = Field(description="List of slide content")

# Tool Schema Models

class GenerateContentSchema(BaseModel):
    """Schema for generate_content tool."""
    request: PresentationRequest = Field(description="Presentation request parameters")

class CreatePowerPointSchema(BaseModel):
    """Schema for create_powerpoint tool."""
    slide_content: SlideContent = Field(description="Content for the slides")
    output_path: str = Field(description="Path where the PowerPoint file should be saved")

class ConvertToImagesSchema(BaseModel):
    """Schema for convert_to_images tool."""
    ppt_path: str = Field(description="Path to the PowerPoint file")
    temp_dir: str = Field(description="Directory to store temporary image files")

class GenerateVideoSchema(BaseModel):
    """Schema for generate_video tool."""
    images_dir: str = Field(description="Directory containing the slide images")
    output_path: str = Field(description="Path where the video file should be saved")

# Content Generation Tools

class ChartElement(SlideElement):
    """Handles charts and graphs."""
    def __init__(self, x: int, y: int, width: int, height: int,
                 chart_type: str, data: Dict[str, List[float]],
                 labels: List[str], colors: List[str]) -> None:
        super().__init__(x, y, width, height)
        self.chart_type = chart_type
        self.data = data
        self.labels = labels
        self.colors = colors
    
    def draw(self, context: cairo.Context) -> None:
        context.save()
        if self.chart_type == "bar":
            self._draw_bar_chart(context)
        elif self.chart_type == "line":
            # Placeholder for line chart implementation
            raise NotImplementedError("Line chart drawing not implemented.")
        elif self.chart_type == "pie":
            # Placeholder for pie chart implementation
            raise NotImplementedError("Pie chart drawing not implemented.")
        context.restore()
    
    def _draw_bar_chart(self, context: cairo.Context) -> None:
        """Draw a simple bar chart."""
        num_series = len(self.data)
        if num_series == 0:
            return
        # Calculate bar width with some spacing (1.5 factor)
        bar_width = self.width / (num_series * 1.5)
        # Find the maximum value among all series
        max_value = max(max(series) for series in self.data.values())
        
        for i, (series_name, values) in enumerate(self.data.items()):
            for j, value in enumerate(values):
                height = (value / max_value) * self.height
                x = self.x + (i + j * 1.5) * bar_width
                y = self.y + self.height - height
                
                # Parse color for the series; use different variable names to avoid shadowing
                r, g, b = [int(self.colors[i][k:k+2], 16)/255 for k in (1, 3, 5)]
                context.set_source_rgb(r, g, b)
                context.rectangle(x, y, bar_width, height)
                context.fill()

def process_slide_shape(shape: Any) -> Optional[SlideElement]:
    """Convert a PowerPoint shape to a SlideElement."""
    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        # Validate that text_frame and paragraphs exist
        font_size = 12
        font_family = "Arial"
        if shape.text_frame and shape.text_frame.paragraphs:
            para = shape.text_frame.paragraphs[0]
            if para.font.size:
                font_size = para.font.size.pt
            if para.font.name:
                font_family = para.font.name
        return TextElement(
            x=shape.left,
            y=shape.top,
            width=shape.width,
            height=shape.height,
            text=shape.text,
            font_size=font_size,
            font_family=font_family,
            color="#000000"  # Default color
        )
    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return ShapeElement(
            x=shape.left,
            y=shape.top,
            width=shape.width,
            height=shape.height,
            shape_type="rectangle",  # Default shape type
            fill_color="#FFFFFF"  # Default color
        )
    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        # Extract chart data if available
        chart_data: Dict[str, List[float]] = {}
        labels: List[str] = []
        try:
            for series in shape.chart.series:
                chart_data[series.name] = list(series.values)
                labels.extend(series.categories)
        except Exception as e:
            logger.error(f"Error processing chart shape: {e}")
        return ChartElement(
            x=shape.left,
            y=shape.top,
            width=shape.width,
            height=shape.height,
            chart_type="bar",  # Default chart type
            data=chart_data,
            labels=labels,
            colors=["#0066CC", "#FF9900", "#33CC33"]  # Default colors
        )
    return None

def export_slide_to_image(slide: Any, output_path: str,
                          config: ImageGenerationConfig,
                          theme: SlideTheme) -> bool:
    """
    Convert a slide to an image using cairo.
    
    Parameters:
        slide: The slide object from the PowerPoint presentation.
        output_path: Path where the PNG image will be saved.
        config: Image generation configuration.
        theme: Slide theme configuration.
    
    Returns:
        True if image export succeeds; otherwise, False.
    """
    try:
        # Create cairo surface
        surface = cairo.ImageSurface(cairo.FORMAT_ARGB32, config.width, config.height)
        context = cairo.Context(surface)
        
        # Set background color
        r, g, b = [int(theme.background_color[i:i+2], 16)/255 for i in (1, 3, 5)]
        context.set_source_rgb(r, g, b)
        context.paint()
        
        # Process and draw each shape in the slide
        elements: List[SlideElement] = []
        for shape in slide.shapes:
            element = process_slide_shape(shape)
            if element:
                elements.append(element)
        
        for element in elements:
            element.draw(context)
        
        # Save the image as PNG
        surface.write_to_png(output_path)
        return True
        
    except Exception as e:
        logger.error(f"Error converting slide to image: {e}")
        return False

def create_transition_frames(image_path: str, output_dir: str, 
                             slide_index: int, num_frames: int) -> None:
    """
    Create transition frames with fade-in effect.
    
    Parameters:
        image_path: Path to the base slide image.
        output_dir: Directory where transition frames will be saved.
        slide_index: Index of the current slide.
        num_frames: Number of transition frames to generate.
    """
    try:
        base_image = Image.open(image_path)
        for frame in range(num_frames):
            # Calculate transition progress (0 to 1)
            progress = frame / (num_frames - 1) if num_frames > 1 else 1
            frame_image = base_image.copy()
            # Apply fade-in effect using brightness enhancement
            enhancer = ImageEnhance.Brightness(frame_image)
            frame_image = enhancer.enhance(progress)
            # Save transition frame image
            frame_path = os.path.join(
                output_dir,
                f"slide_{slide_index:03d}_transition_{frame:02d}.jpg"
            )
            frame_image.save(frame_path, quality=95)
    except Exception as e:
        logger.error(f"Error creating transition frames: {e}")

@tool("generate_content", args_schema=GenerateContentSchema)
def generate_content(request: PresentationRequest) -> SlideContent:
    """Generates structured content for presentation slides."""
    logger.info(f"Generating content for {request.num_slides} slides about {request.main_topic}")
    
    slides: List[SlideData] = []
    
    # Introduction slide content and bullet points
    intro_bullets = [
        f"Understanding the impact and potential of {request.main_topic}",
        "Key developments and current state of the field",
        "Future prospects and opportunities",
        "Implementation considerations and best practices"
    ]
    
    intro_content = f"""
Welcome to this comprehensive overview of {request.main_topic}. This presentation will explore the current landscape, future developments, and practical implications of this transformative field.

We will examine key aspects including historical context, current applications, challenges, and opportunities. Our discussion will cover both theoretical foundations and practical implementation considerations."""
    
    slides.append(SlideData(
        topic=request.main_topic,
        content=intro_content,
        bullet_points=intro_bullets
    ))
    
    # Define main content topics
    main_topics = [
        "Background and Context",
        "Current State",
        "Future Developments",
        "Challenges and Opportunities",
        "Best Practices",
        "Case Studies",
        "Implementation Strategies",
        "Key Considerations"
    ]
    
    def generate_bullet_points(topic: str, main_topic: str, num_points: int = 3) -> List[str]:
        """Generate bullet points based on topic and main subject."""
        bullet_points_map = {
            "Background and Context": [
                f"Historical development of {main_topic} from early concepts to present day",
                f"Key milestones and breakthroughs in {main_topic}",
                f"Evolution of technologies and methodologies in {main_topic}"
            ],
            "Current State": [
                f"Latest advancements and breakthroughs in {main_topic}",
                f"Major players and organizations driving innovation in {main_topic}",
                f"Current challenges and limitations facing {main_topic}"
            ],
            "Future Developments": [
                f"Predicted trends and innovations in {main_topic} for the next decade",
                f"Potential breakthrough technologies that could transform {main_topic}",
                f"Expected impact on industries and society"
            ],
            "Challenges and Opportunities": [
                f"Technical challenges in advancing {main_topic}",
                f"Ethical considerations and regulatory frameworks",
                f"Opportunities for innovation and growth in {main_topic}"
            ],
            "Best Practices": [
                f"Industry standards and guidelines for {main_topic}",
                f"Security and safety considerations in implementing {main_topic}",
                f"Quality assurance and testing methodologies"
            ],
            "Case Studies": [
                f"Successful implementations of {main_topic} in various industries",
                f"Lessons learned from early adopters and pioneers",
                f"Real-world impact and results from {main_topic} applications"
            ],
            "Implementation Strategies": [
                f"Step-by-step approach to adopting {main_topic}",
                f"Resource requirements and infrastructure considerations",
                f"Risk mitigation and change management strategies"
            ],
            "Key Considerations": [
                f"Critical factors for successful implementation of {main_topic}",
                f"Cost-benefit analysis and ROI considerations",
                f"Long-term sustainability and maintenance aspects"
            ]
        }
        return bullet_points_map.get(topic, [
            f"Key aspect {i+1} of {topic} in {main_topic}"
            for i in range(num_points)
        ])
    
    def generate_slide_content(topic: str, main_topic: str) -> str:
        """Generate detailed slide content based on topic and main subject."""
        content_map = {
            "Background and Context": f"""
The evolution of {main_topic} represents a fascinating journey of technological advancement and innovation. Starting from its conceptual foundations, this field has undergone significant transformation through various developmental stages.

Key historical developments have shaped our current understanding and application of {main_topic}. These developments include breakthrough discoveries, technological innovations, and paradigm shifts in approaching complex challenges.

Understanding this historical context is crucial for appreciating the current state and future potential of {main_topic}.""",
            "Current State": f"""
{main_topic} is currently experiencing rapid advancement and widespread adoption across various sectors. The field has reached a critical point where theoretical concepts are being successfully translated into practical applications.

Modern implementations of {main_topic} demonstrate sophisticated capabilities and are solving increasingly complex problems. The current landscape is characterized by intense research activity and innovative applications.

Industry leaders and research institutions are continuously pushing the boundaries of what's possible with {main_topic}, leading to new breakthroughs and applications.""",
            "Future Developments": f"""
The future of {main_topic} holds immense promise and potential for transformative impact across industries. Emerging trends and technological advances suggest significant developments on the horizon.

Anticipated breakthroughs in {main_topic} are expected to revolutionize current practices and enable new possibilities. These developments will likely reshape how we approach challenges and create solutions.

The convergence of {main_topic} with other emerging technologies is expected to create new opportunities and applications.""",
            "Challenges and Opportunities": f"""
While {main_topic} presents tremendous opportunities, it also faces significant challenges that need to be addressed. Technical limitations, resource constraints, and implementation complexities are key areas requiring attention.

However, these challenges also present opportunities for innovation and improvement. The field is ripe for breakthrough solutions and novel approaches to overcome current limitations.

Balancing the potential benefits with practical considerations is crucial for successful advancement in {main_topic}.""",
            "Best Practices": f"""
Implementing {main_topic} effectively requires adherence to established best practices and industry standards. These guidelines ensure optimal performance, reliability, and safety in applications.

Key considerations include proper planning, thorough testing, and continuous monitoring of systems. Security measures and quality assurance processes play vital roles in successful implementations.

Regular updates and maintenance procedures help maintain system effectiveness and adapt to evolving requirements.""",
            "Case Studies": f"""
Real-world implementations of {main_topic} provide valuable insights into its practical applications and impact. Successful case studies demonstrate the tangible benefits and lessons learned from various deployments.

Organizations across different sectors have shown how {main_topic} can be effectively utilized to solve complex problems and create value. These examples provide valuable guidance for future implementations.

Learning from both successes and challenges in these cases helps improve future applications of {main_topic}.""",
            "Implementation Strategies": f"""
Successfully implementing {main_topic} requires a well-planned and systematic approach. Key strategies include thorough assessment of requirements, careful planning of resources, and staged implementation.

Important considerations include infrastructure requirements, team capabilities, and change management processes. Risk assessment and mitigation strategies play crucial roles in successful deployment.

Regular evaluation and adjustment of implementation strategies ensure optimal results and long-term success.""",
            "Key Considerations": f"""
Several critical factors must be considered when working with {main_topic}. These include technical requirements, resource allocation, and long-term sustainability.

Cost-benefit analysis and return on investment calculations help determine the viability of {main_topic} implementations. Understanding and planning for maintenance and upgrade requirements is essential.

Successful integration of {main_topic} depends on careful consideration of these factors and appropriate planning."""
        }
        return content_map.get(topic, f"Detailed exploration of {topic} in the context of {main_topic}")
    
    # Generate main slides based on the topics
    for i in range(1, request.num_slides):
        topic_idx = i % len(main_topics)
        current_topic = main_topics[topic_idx]
        slide_topic = f"{current_topic} in {request.main_topic}"
        
        content = generate_slide_content(current_topic, request.main_topic)
        bullets = generate_bullet_points(current_topic, request.main_topic)
        
        slides.append(SlideData(
            topic=slide_topic,
            content=content,
            bullet_points=bullets
        ))
    
    logger.info(f"Generated content for {len(slides)} slides")
    return SlideContent(slides=slides)

@tool("create_powerpoint", args_schema=CreatePowerPointSchema)
def create_powerpoint(slide_content: SlideContent, output_path: str) -> str:
    """Creates a PowerPoint presentation from structured slide content."""
    logger.info(f"Creating PowerPoint with {len(slide_content.slides)} slides")
    prs = Presentation()
    
    try:
        for idx, slide_data in enumerate(slide_content.slides):
            logger.debug(f"Creating slide {idx+1}: {slide_data.topic}")
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.topic
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(32)
            
            # Add content (combining main content and bullet points)
            content_placeholder = slide.placeholders[1]
            content_text = slide_data.content
            if slide_data.bullet_points:
                content_text += "\n\n"
                for bullet in slide_data.bullet_points:
                    content_text += f"• {bullet}\n"
            content_placeholder.text = content_text
            
            # Format content text
            for paragraph in content_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(16)
        
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        
        prs.save(output_path)
        logger.info(f"PowerPoint saved to {output_path}")
        return output_path
    
    except Exception as e:
        logger.error(f"Error creating PowerPoint: {e}")
        raise

@tool("convert_to_images", args_schema=ConvertToImagesSchema)
def convert_to_images(ppt_path: str, temp_dir: str) -> str:
    """Converts PowerPoint slides to images with enhanced rendering."""
    logger.info(f"Converting PowerPoint to images: {ppt_path} -> {temp_dir}")
    
    # Remove temporary directory if it exists and then recreate it
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
    Path(temp_dir).mkdir(parents=True, exist_ok=True)
    
    try:
        # Configuration and theme
        config = ImageGenerationConfig()
        theme = SlideTheme()
        
        # Open the presentation
        presentation = Presentation(ppt_path)
        
        for slide_index, slide in enumerate(presentation.slides, 1):
            # Define static image path
            static_image_path = os.path.join(temp_dir, f"slide_{slide_index:03d}.jpg")
            success = export_slide_to_image(slide, static_image_path, config, theme)
            
            if success and config.enable_animations:
                # Create transition frames for the slide
                create_transition_frames(
                    static_image_path,
                    temp_dir,
                    slide_index,
                    config.transition_frames
                )
        return temp_dir
        
    except Exception as e:
        logger.error(f"Error converting to images: {e}")
        raise

@tool("generate_video", args_schema=GenerateVideoSchema)
def generate_video(images_dir: str, output_path: str) -> str:
    """Creates a video from slide images with transitions."""
    logger.info(f"Generating video from images in {images_dir}")
    
    images = sorted([
        img for img in os.listdir(images_dir) 
        if img.lower().endswith((".jpg", ".jpeg", ".png"))
    ])
    
    if not images:
        raise Exception("No images found in the specified directory")
    
    try:
        first_image_path = os.path.join(images_dir, images[0])
        first_image = cv2.imread(first_image_path)
        if first_image is None:
            raise Exception(f"Could not read the first image: {first_image_path}")
        height, width, layers = first_image.shape
        
        video = cv2.VideoWriter(
            output_path,
            cv2.VideoWriter_fourcc(*'mp4v'),
            30,
            (width, height)
        )
        
        for image_name in images:
            image_path = os.path.join(images_dir, image_name)
            frame = cv2.imread(image_path)
            if frame is None:
                logger.error(f"Skipping invalid frame: {image_path}")
                continue
            # If the image is a transition frame, write it once;
            # Otherwise, hold the static slide for 3 seconds (90 frames at 30 fps)
            if "transition" in image_name:
                video.write(frame)
            else:
                for _ in range(90):
                    video.write(frame)
        
        video.release()
        logger.info(f"Video generated: {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"Error generating video: {e}")
        raise

def create_agents() -> Tuple[Any, Any, Any, Any]:
    """Creates specialized agents for different tasks."""
    logger.info("Creating specialized agents")
    
    google_api_key = os.getenv("GOOGLE_API_KEY")
    if not google_api_key:
        logger.error("GOOGLE_API_KEY not found in environment variables.")
        raise Exception("Missing GOOGLE_API_KEY")
    
    model = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        google_api_key=google_api_key
    )
    
    # Content Generation Agent
    content_agent = create_react_agent(
        model=model,
        tools=[generate_content],
        name="content_expert",
        prompt=(
            "You are an expert content creator specializing in presentations. "
            "Generate informative, engaging, and well-structured content that flows naturally. "
            "Focus on creating content that is both educational and engaging."
        )
    )
    
    # PowerPoint Creation Agent
    ppt_agent = create_react_agent(
        model=model,
        tools=[create_powerpoint],
        name="ppt_expert",
        prompt=(
            "You are a PowerPoint expert. Create well-structured presentations "
            "that effectively communicate the content while maintaining visual appeal."
        )
    )
    
    # Image Conversion Agent
    image_agent = create_react_agent(
        model=model,
        tools=[convert_to_images],
        name="image_expert",
        prompt=(
            "You are an expert in converting presentations to high-quality images. "
            "Ensure proper resolution and smooth transitions between slides."
        )
    )
    
    # Video Creation Agent
    video_agent = create_react_agent(
        model=model,
        tools=[generate_video],
        name="video_expert",
        prompt=(
            "You are an expert in creating video presentations. "
            "Create professional videos with smooth transitions and appropriate timing."
        )
    )
    
    return content_agent, ppt_agent, image_agent, video_agent

def create_presentation_workflow() -> Any:
    """Creates and configures the presentation creation workflow."""
    logger.info("Creating presentation workflow")
    
    content_agent, ppt_agent, image_agent, video_agent = create_agents()
    
    workflow = create_supervisor(
        [content_agent, ppt_agent, image_agent, video_agent],
        model=ChatGoogleGenerativeAI(
            model="gemini-2.0-flash",
            google_api_key=os.getenv("GOOGLE_API_KEY")
        ),
        prompt=(
            "You are a project supervisor coordinating a team of experts to create video presentations. "
            "Your team includes:\n"
            "- content_expert: Creates engaging presentation content\n"
            "- ppt_expert: Builds professional PowerPoint presentations\n"
            "- image_expert: Converts presentations to high-quality images\n"
            "- video_expert: Creates polished video presentations\n"
            "Guide the team to create high-quality video presentations efficiently."
        )
    )
    
    return workflow.compile()

def create_presentation_video(topic: str, num_slides: int, target_audience: str = "general") -> str:
    """
    Main function to create a video presentation.
    
    Parameters:
        topic: The main topic of the presentation.
        num_slides: Number of slides in the presentation.
        target_audience: Intended audience for the presentation.
    
    Returns:
        A success message with the path to the created video.
    """
    logger.info(f"Creating presentation video about {topic} with {num_slides} slides")
    
    current_dir = os.path.abspath(os.path.dirname(__file__))
    temp_ppt_path = os.path.join(current_dir, "temp_presentation.pptx")
    temp_slides_dir = os.path.join(current_dir, "temp_slides")
    video_output_path = os.path.join(current_dir, "presentation_video.mp4")
    
    try:
        # Create presentation request
        request = PresentationRequest(
            num_slides=num_slides,
            main_topic=topic,
            target_audience=target_audience
        )
        
        # Step 1: Generate content
        logger.info("Generating content")
        content_result = generate_content.invoke({"request": request})
        
        # Step 2: Create PowerPoint
        logger.info("Creating PowerPoint")
        ppt_path = create_powerpoint.invoke({
            "slide_content": content_result,
            "output_path": temp_ppt_path
        })
        
        # Step 3: Convert PowerPoint slides to images
        logger.info("Converting slides to images")
        images_dir = convert_to_images.invoke({
            "ppt_path": ppt_path,
            "temp_dir": temp_slides_dir
        })
        
        # Step 4: Generate video from images
        logger.info("Generating video")
        video_path = generate_video.invoke({
            "images_dir": images_dir,
            "output_path": video_output_path
        })
        
        if os.path.exists(video_path):
            video_size = os.path.getsize(video_path) / (1024 * 1024)  # Size in MB
            logger.info(f"Video created successfully. Size: {video_size:.2f} MB")
            return f"Video created successfully at: {video_path}"
        else:
            raise Exception("Video file not found after creation")
            
    except Exception as e:
        logger.error(f"Error creating presentation: {e}")
        raise
    finally:
        # Cleanup temporary files
        logger.info("Cleaning up temporary files")
        try:
            if os.path.exists(temp_ppt_path):
                os.remove(temp_ppt_path)
            if os.path.exists(temp_slides_dir):
                shutil.rmtree(temp_slides_dir)
        except Exception as e:
            logger.warning(f"Error during cleanup: {e}")

if __name__ == "__main__":
    try:
        # Example usage
        result = create_presentation_video(
            topic="The Advancement in LLM's",
            num_slides=5,
            target_audience="technical"
        )
        print(result)
    except Exception as e:
        print(f"Error: {e}")
